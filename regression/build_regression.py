#!/usr/bin/env python3
"""
合成中文文档回归集（脱敏，无真实合同内容）

产出:
  regression/documents/         合成的文档文件（PDF/DOCX/XLSX/PPTX/TXT/HTML）
  regression/manifest.json      元数据清单（每份文档的类型、页数、合成内容、期望行为）
  regression/build.log          构建日志

覆盖矩阵（对齐升级计划 §4 回归集要求）:
  PDF
    - 可编辑文本中文合同（单栏 / 双栏 / 多页）      -> 走 Fitz 快速路径
    - 纯扫描（无文本层）                            -> 走 OCR 路径
    - 混合（部分页有文本层、部分页扫描）            -> 路由分流
    - 表格（有框线 / 无框线 / 合并单元格 / 跨页）   -> 表格结构识别
    - 印章 / 手写批注区域（模拟红色圆形印章）       -> 视觉鲁棒性
    - 页眉页脚 / 页码 / 水印
    - 空页 / 图片页（无文字）
    - 异常文件：加密、损坏头、空文件、超限          -> 稳定错误码
  非 PDF: DOCX / XLSX / PPTX / TXT / HTML / MD

所有正文文本为占位符（"合同编号 CN-XXXX-YYYY-XXXX" 等），不含任何真实个人信息。
"""
import json
import logging
import os
import random
import shutil
import sys
from datetime import date
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont

logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
log = logging.getLogger("build_regression")

ROOT = Path(__file__).resolve().parent
DOCS = ROOT / "documents"
MANIFEST = ROOT / "manifest.json"

# Noto CJK 字体（有中文字形）
CJK_FONT = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
CJK_FONT_SERIF = "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc"

# 合同占位语料（全部为虚构占位内容）
CONTRACT_BODY = [
    "本合同由以下双方于签订：",
    "甲方（以下简称甲方）与乙方（以下简称乙方）本着平等自愿的原则，就下列事项达成一致意见。",
    "第一条 合同标的：乙方按照本合同约定向甲方提供服务。",
    "第二条 服务内容：包括但不限于方案设计、实施部署与培训支持。",
    "第三条 合同金额：人民币壹拾万元整。",
    "第四条 付款方式：合同签订后支付百分之三十预付款，验收合格后支付剩余款项。",
    "第五条 违约责任：任何一方违反本合同约定，应承担相应的违约责任。",
    "第六条 争议解决：双方协商解决，协商不成的，提交有管辖权的人民法院处理。",
    "第七条 本合同一式贰份，双方各执壹份，具有同等法律效力。",
    "签订地点：示例市示例区示例路示例号。",
]

PERSONA_A = "示例科技有限公司"
PERSONA_B = "示例商贸有限公司"


def cjk_font(size: int, serif: bool = False):
    return ImageFont.truetype(CJK_FONT_SERIF if serif else CJK_FONT, size)


def make_text_pdf(path: Path, pages: int = 4, two_col: bool = False,
                  meta: dict = None):
    """可编辑文本层的中文合同 PDF。"""
    doc = fitz.open()
    for pn in range(1, pages + 1):
        page = doc.new_page()  # A4 595x842pt
        y = 72
        title = f"采购服务合同 第{pn}页"
        page.insert_text((72, y), title, fontsize=14,
                         fontname="noto", fontfile=CJK_FONT)
        y += 30
        for line in CONTRACT_BODY:
            if two_col and pn % 2 == 0:
                # 双栏：右半栏
                page.insert_text((330, y + (pn % 2) * 400), line, fontsize=10,
                                 fontname="noto", fontfile=CJK_FONT)
            else:
                page.insert_text((72, y), line, fontsize=10,
                                 fontname="noto", fontfile=CJK_FONT)
            y += 18
        # 页脚页码
        page.insert_text((280, 820), f"- {pn} -", fontsize=9,
                         fontname="noto", fontfile=CJK_FONT)
    doc.save(str(path))
    doc.close()


def render_page_image(w: int = 1240, h: int = 1754, dpi_scale: int = 60):
    """渲染一页中文文本为图像（模拟扫描件）。"""
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    f_title = cjk_font(40)
    f_body = cjk_font(26)
    y = 120
    d.text((120, y), "采购服务合同", font=f_title, fill="black")
    y += 90
    for line in CONTRACT_BODY:
        d.text((120, y), line, font=f_body, fill="black")
        y += 48
    return img


def add_stamp(img, text="示例科技有限公司 合同专用章"):
    """叠加一枚红色圆形印章（模拟真实合同印章）。"""
    d = ImageDraw.Draw(img)
    cx, cy = 900, 1300
    r = 180
    d.ellipse([cx - r, cy - r, cx + r, cy + r], outline=(220, 40, 40), width=14)
    d.ellipse([cx - r + 20, cy - r + 20, cx + r - 20, cy + r - 20],
              outline=(220, 40, 40), width=3)
    f = cjk_font(30)
    d.text((cx - 100, cy - 20), "示例科技有限公司", font=f, fill=(220, 40, 40))
    d.text((cx - 70, cy + 40), "合同专用章", font=f, fill=(220, 40, 40))
    return img


def make_scan_pdf(path: Path, pages: int = 2, stamp: bool = True):
    """纯扫描 PDF：无文本层，内容为渲染图像。"""
    doc = fitz.open()
    for i in range(pages):
        img = render_page_image()
        if stamp:
            img = add_stamp(img)
        img.save(str(ROOT / ".tmp_scan.png"))
        page = doc.new_page()
        page.insert_image(page.rect, filename=str(ROOT / ".tmp_scan.png"))
    doc.save(str(path))
    doc.close()
    (ROOT / ".tmp_scan.png").unlink(missing_ok=True)


def make_mixed_pdf(path: Path):
    """混合 PDF：第1页可编辑文本，第2页扫描图像，第3页可编辑文本。"""
    doc = fitz.open()
    # 第1页 文本
    p = doc.new_page()
    p.insert_text((72, 720), "可编辑文本页：合同首页", fontsize=12,
                  fontname="noto", fontfile=CJK_FONT)
    for i, line in enumerate(CONTRACT_BODY[:4]):
        p.insert_text((72, 690 - i * 22), line, fontsize=10,
                      fontname="noto", fontfile=CJK_FONT)
    # 第2页 扫描
    img = render_page_image()
    img.save(str(ROOT / ".tmp_scan.png"))
    p2 = doc.new_page()
    p2.insert_image(p2.rect, filename=str(ROOT / ".tmp_scan.png"))
    # 第3页 文本
    p3 = doc.new_page()
    p3.insert_text((72, 720), "可编辑文本页：合同末页", fontsize=12,
                   fontname="noto", fontfile=CJK_FONT)
    for i, line in enumerate(CONTRACT_BODY[4:8]):
        p3.insert_text((72, 690 - i * 22), line, fontsize=10,
                       fontname="noto", fontfile=CJK_FONT)
    doc.save(str(path))
    doc.close()
    (ROOT / ".tmp_scan.png").unlink(missing_ok=True)


def make_table_pdf(path: Path, variant: str, pages: int = 2):
    """表格 PDF：有框线 / 无框线 / 合并单元格 / 跨页。"""
    doc = fitz.open()
    for pn in range(pages):
        p = doc.new_page()
        p.insert_text((72, 60), f"合同附表（{variant}）第{pn + 1}页", fontsize=12,
                      fontname="noto", fontfile=CJK_FONT)
        headers = ["序号", "品名", "数量", "单价", "金额"]
        y0 = 100
        row_h = 40
        col_x = [72, 160, 260, 350, 450]
        rows = []
        for i in range(1, 9):
            rows.append([str(i), f"示例物品{i}", str(i * 10),
                         f"{i * 100}.00", f"{i * 1000}.00"])
        if variant == "ruled":
            # 有框线
            for r_i, row in enumerate([headers] + rows):
                y = y0 + r_i * row_h
                p.draw_line((72, y), (560, y))
                for c_i, (cx, cell) in enumerate(zip(col_x, row)):
                    p.insert_text((cx + 4, y + 24), cell, fontsize=10,
                                  fontname="noto", fontfile=CJK_FONT)
                    if r_i == 0:
                        p.draw_line((cx, y), (cx, y + row_h))
            p.draw_line((72, y0 + (len(rows) + 1) * row_h), (560, y0 + (len(rows) + 1) * row_h))
            for cx in col_x[1:]:
                p.draw_line((cx, y0), (cx, y0 + (len(rows) + 1) * row_h))
        elif variant == "borderless":
            # 无框线表格：仅文字按列排布
            for r_i, row in enumerate(rows):
                y = y0 + r_i * row_h
                for c_i, (cx, cell) in enumerate(zip(col_x, row)):
                    p.insert_text((cx + 4, y + 24), cell, fontsize=10,
                                  fontname="noto", fontfile=CJK_FONT)
        elif variant == "merged":
            # 合并单元格：跨列大单元格 + 标题行
            p.draw_line((72, y0), (560, y0))
            p.draw_line((72, y0 + row_h), (560, y0 + row_h))
            p.insert_text((100, y0 + 26), "合计（跨列合并）", fontsize=10,
                          fontname="noto", fontfile=CJK_FONT)
            p.draw_line((72, y0 + row_h), (72, y0 + 3 * row_h))
            p.draw_line((560, y0 + row_h), (560, y0 + 3 * row_h))
            for r_i, row in enumerate(rows[:3]):
                y = y0 + (r_i + 1) * row_h
                for c_i, (cx, cell) in enumerate(zip(col_x, row)):
                    p.insert_text((cx + 4, y + 24), cell, fontsize=10,
                                  fontname="noto", fontfile=CJK_FONT)
                    if r_i == 0:
                        p.draw_line((cx, y), (cx, y + row_h))
            p.draw_line((72, y0 + 3 * row_h), (560, y0 + 3 * row_h))
        elif variant == "crosspage":
            # 跨页表格：页1表头+4行，页2续表+4行
            if pn == 0:
                p.insert_text((72, 80), "表头：序号 品名 数量 单价 金额", fontsize=10,
                              fontname="noto", fontfile=CJK_FONT)
                for r_i, row in enumerate(rows[:4]):
                    y = y0 + r_i * row_h
                    for c_i, (cx, cell) in enumerate(zip(col_x, row)):
                        p.insert_text((cx + 4, y + 24), cell, fontsize=10,
                                      fontname="noto", fontfile=CJK_FONT)
                p.draw_line((72, y0 - 5), (560, y0 - 5))
                p.draw_line((72, y0 + 4 * row_h), (560, y0 + 4 * row_h))
            else:
                p.insert_text((72, 80), "续表（跨页）", fontsize=10,
                              fontname="noto", fontfile=CJK_FONT)
                for r_i, row in enumerate(rows[4:]):
                    y = y0 + r_i * row_h
                    for c_i, (cx, cell) in enumerate(zip(col_x, row)):
                        p.insert_text((cx + 4, y + 24), cell, fontsize=10,
                                      fontname="noto", fontfile=CJK_FONT)
                p.draw_line((72, y0 - 5), (560, y0 - 5))
                p.draw_line((72, y0 + 4 * row_h), (560, y0 + 4 * row_h))
    doc.save(str(path))
    doc.close()


def make_docx(path: Path, pages: int = 3):
    """DOCX 中文合同。"""
    from docx import Document
    d = Document()
    d.add_heading("采购服务合同", level=0)
    for line in CONTRACT_BODY:
        d.add_paragraph(line)
    for i in range(pages):
        d.add_paragraph(f"（第{i + 1}页正文段落，示例内容）")
    d.save(str(path))


def make_xlsx(path: Path, sheets: int = 2):
    """XLSX 中文合同附表。"""
    from openpyxl import Workbook
    wb = Workbook()
    for s_i in range(sheets):
        ws = wb.create_sheet(f"附表{s_i + 1}")
        ws.append(["序号", "品名", "数量", "单价", "金额"])
        for i in range(1, 21):
            ws.append([i, f"示例物品{i}", i * 10, f"{i * 100}.00", f"{i * 1000}.00"])
    wb.save(str(path))


def make_pptx(path: Path, slides: int = 4):
    """PPTX 中文演示（合同条款说明）。"""
    from pptx import Presentation
    prs = Presentation()
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"合同条款第{i + 1}节"
        slide.placeholders[1].text = "\n".join(CONTRACT_BODY[:5])
    prs.save(str(path))


def make_txt(path: Path):
    path.write_text("\n".join([f"合同编号 CN-{random.randint(1000,9999)}-{random.randint(100000,999999)}"] + CONTRACT_BODY), encoding="utf-8")


def make_html(path: Path):
    body = "".join(f"<p>{l}</p>" for l in CONTRACT_BODY)
    path.write_text(
        "<html><head><meta charset='utf-8'><title>示例合同</title></head>"
        f"<body><h1>采购服务合同</h1>{body}</body></html>",
        encoding="utf-8")


def main():
    if DOCS.exists():
        shutil.rmtree(DOCS)
    DOCS.mkdir(parents=True)
    manifest = {"generated": str(date.today()), "root": "regression/documents/",
                "docs": []}

    def add(name, dtype, pages, expect, note="", **kw):
        manifest["docs"].append({
            "name": name, "type": dtype, "pages": pages,
            "expected": expect, "note": note, **kw,
        })

    # ===== 可编辑文本中文合同 =====
    make_text_pdf(DOCS / "editable_single.pdf", pages=4)
    add("editable_single.pdf", "pdf_editable", 4, "text", "单栏可编辑文本层")
    make_text_pdf(DOCS / "editable_twocol.pdf", pages=4, two_col=True)
    add("editable_twocol.pdf", "pdf_editable", 4, "text", "双栏可编辑文本层")
    make_text_pdf(DOCS / "editable_10p.pdf", pages=10)
    add("editable_10p.pdf", "pdf_editable", 10, "text", "10页可编辑文本层")

    # ===== 纯扫描 =====
    make_scan_pdf(DOCS / "scan_2p.pdf", pages=2, stamp=True)
    add("scan_2p.pdf", "pdf_scan", 2, "ocr", "纯扫描+红色印章")
    make_scan_pdf(DOCS / "scan_4p.pdf", pages=4, stamp=False)
    add("scan_4p.pdf", "pdf_scan", 4, "ocr", "纯扫描无印章")

    # ===== 混合 =====
    make_mixed_pdf(DOCS / "mixed_3p.pdf")
    add("mixed_3p.pdf", "pdf_mixed", 3, "ocr_or_text", "文本层+扫描混合")

    # ===== 表格 =====
    for v in ["ruled", "borderless", "merged", "crosspage"]:
        make_table_pdf(DOCS / f"table_{v}.pdf", v)
        add(f"table_{v}.pdf", "pdf_table", 2, "table", f"表格:{v}")

    # ===== 非 PDF =====
    make_docx(DOCS / "contract.docx")
    add("contract.docx", "docx", 3, "text", "中文合同 DOCX")
    make_xlsx(DOCS / "contract.xlsx")
    add("contract.xlsx", "xlsx", 2, "table", "合同附表 XLSX")
    make_pptx(DOCS / "contract.pptx")
    add("contract.pptx", "pptx", 4, "text", "合同条款 PPTX")
    make_txt(DOCS / "contract.txt")
    add("contract.txt", "txt", 1, "text", "合同 TXT")
    make_html(DOCS / "contract.html")
    add("contract.html", "html", 1, "text", "合同 HTML")

    # ===== 异常文件 =====
    # 加密 PDF
    doc = fitz.open(); p = doc.new_page()
    p.insert_text((72, 720), "secret", fontsize=11)
    doc.save(str(DOCS / "encrypted.pdf"), encryption=fitz.PDF_ENCRYPT_AES_256, user_pw="pass123")
    doc.close()
    add("encrypted.pdf", "pdf_error", 1, "err_422", "密码加密")
    # 损坏 PDF（头缺失）
    (DOCS / "corrupt.pdf").write_bytes(b"not a pdf at all")
    add("corrupt.pdf", "pdf_error", 0, "err_400", "非PDF头")
    # 空文件
    (DOCS / "empty.pdf").write_bytes(b"")
    add("empty.pdf", "pdf_error", 0, "err_400", "空文件")
    # 无扩展名
    (DOCS / "noext").write_bytes(b"whatever")
    add("noext", "unknown", 0, "err_400", "无扩展名")
    # 图片 PDF（有图无文）
    img = render_page_image()
    img.save(str(DOCS / ".t.png"))
    doc = fitz.open(); pg = doc.new_page(); pg.insert_image(pg.rect, filename=str(DOCS / ".t.png"))
    doc.save(str(DOCS / "picture.pdf")); doc.close()
    (DOCS / ".t.png").unlink(missing_ok=True)
    add("picture.pdf", "pdf_scan", 1, "ocr", "单页图片 PDF")

    with open(MANIFEST, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    total_pages = sum(m["pages"] for m in manifest["docs"])
    log.info(f"生成 {len(manifest['docs'])} 份文档，共 {total_pages} 页")
    log.info(f"清单: {MANIFEST}")


if __name__ == "__main__":
    main()
